"""Caller-owned preset library works without private backend assets."""
import json
import os
from pathlib import Path
import subprocess
import sys

from pptx import Presentation
import pytest

RENDERER = Path(__file__).resolve().parents[1] / 'tools/pptx_creator.py'


def library(tmp_path):
    pack = tmp_path / 'uploaded-pack'
    (pack / 'master').mkdir(parents=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'Synthetic customer-owned slide'
    prs.save(pack / 'source.pptx')
    (pack / 'master/slide-index.json').write_text(json.dumps({
        'slides': [{'id': 'title', 'source_file': 'source.pptx', 'source_slide_index': 1}],
        'presets': {'example': {'description': 'Synthetic preset', 'product': 'test', 'slides': ['title']}}}))
    return pack


def test_preset_generation_from_uploaded_workspace_library(tmp_path):
    pack = library(tmp_path)
    env = {**os.environ, 'HOME': str(tmp_path), 'LOMA_ISOLATED_WORKER': '1'}
    result = subprocess.run([sys.executable, str(RENDERER), '--library-dir', str(pack),
        'generate', '--preset', 'example', '--client', 'Synthetic', '--output', 'result.pptx'],
        env=env, capture_output=True, text=True, timeout=30)
    assert result.returncode == 0, result.stderr
    deck = Presentation(tmp_path / 'artifacts/result.pptx')
    assert len(deck.slides) == 1
    assert deck.slides[0].shapes.title.text == 'Synthetic customer-owned slide'


@pytest.mark.parametrize('kind', ['outside', 'symlink', 'missing'])
def test_library_must_be_inside_workspace(tmp_path, kind):
    pack = library(tmp_path)
    home = tmp_path / 'home'
    home.mkdir()
    if kind == 'symlink':
        (home / 'link').symlink_to(pack, target_is_directory=True)
        pack = home / 'link'
    elif kind == 'missing':
        pack = home
    result = subprocess.run([sys.executable, str(RENDERER), '--library-dir', str(pack), 'presets'],
        env={**os.environ, 'HOME': str(home), 'LOMA_ISOLATED_WORKER': '1'},
        capture_output=True, text=True, timeout=30)
    assert result.returncode == 2
    assert 'Invalid library' in result.stderr
