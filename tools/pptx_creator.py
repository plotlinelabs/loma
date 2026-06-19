#!/usr/bin/env python3
"""PPTX Presentation Generator — Compose client-specific product decks.

Provides CLI commands for the Loma agent:
  1. pptx_creator.py catalog [--category C] [--product P]  — List available slides
  2. pptx_creator.py layouts                               — List slide layouts for new slides
  3. pptx_creator.py presets                               — List use-case presets
  4. pptx_creator.py assets [--category C]                 — List asset library
  5. pptx_creator.py generate --preset P --client NAME     — Generate from preset
  6. pptx_creator.py compose --spec FILE                   — Generate from JSON deck spec

Usage (called by the agent via Bash):
  python3 tools/pptx_creator.py catalog
  python3 tools/pptx_creator.py catalog --category nudges_tooltips
  python3 tools/pptx_creator.py presets
  python3 tools/pptx_creator.py generate --preset first-call-in-app --client "Example Corp"
  python3 tools/pptx_creator.py compose --spec /path/to/deck_spec.json

Deck spec JSON format:
  {
    "client": "Example Corp",
    "output": "Example-First-Call.pptx",
    "slides": [
      {"action": "use", "slide_id": "title"},
      {"action": "modify", "slide_id": "social_proof", "modifications": {
        "text": [{"shape_name": "Text 0", "value": "500M+"}]
      }},
      {"action": "create", "layout": "feature-showcase", "content": {
        "title": "Custom Feature", "body": "Description", "image": "path/to/image.png"
      }}
    ]
  }
"""

import argparse
import json
import os
import shutil
import sys
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt, Emu

PROJECT_ROOT = Path(__file__).parent.parent
PPT_DIR = PROJECT_ROOT / "ppt"
MASTER_DIR = PPT_DIR / "master"
OUTPUT_DIR = PPT_DIR / "output"
ASSETS_DIR = PPT_DIR / "assets"
SLIDE_INDEX_PATH = MASTER_DIR / "slide-index.json"

# Cache for loaded presentations (avoid re-opening large files)
_prs_cache: dict[str, Presentation] = {}


def _load_slide_index() -> dict:
    if not SLIDE_INDEX_PATH.exists():
        print(json.dumps({"error": f"Slide index not found at {SLIDE_INDEX_PATH}. Run: python3 ppt/build_master.py"}))
        sys.exit(1)
    with open(SLIDE_INDEX_PATH) as f:
        return json.load(f)


def _load_presentation(filename: str) -> Presentation:
    """Load a PPTX file, using cache to avoid re-opening."""
    if filename not in _prs_cache:
        path = PPT_DIR / filename
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {path}")
        _prs_cache[filename] = Presentation(str(path))
    return _prs_cache[filename]


def _resolve_asset_path(asset_ref: str) -> Path:
    """Resolve an asset reference to an absolute path."""
    # Try relative to assets dir
    p = ASSETS_DIR / asset_ref
    if p.exists():
        return p
    # Try relative to ppt dir
    p = PPT_DIR / asset_ref
    if p.exists():
        return p
    # Try as absolute path
    p = Path(asset_ref)
    if p.exists():
        return p
    raise FileNotFoundError(f"Asset not found: {asset_ref}")


# ---------------------------------------------------------------------------
# Slide manipulation functions
# ---------------------------------------------------------------------------

def delete_slide(prs: Presentation, slide_index: int):
    """Delete a slide from the presentation by 0-based index.

    python-pptx has no native delete — manipulate XML directly.
    """
    rId = prs.slides._sldIdLst[slide_index].get("r:id")
    if rId is None:
        # Try without namespace
        rId = prs.slides._sldIdLst[slide_index].attrib.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(sldId)


def reorder_slides(prs: Presentation, new_order: list[int]):
    """Reorder slides. new_order is a list of current 0-based indices in desired order."""
    sldIdLst = prs.slides._sldIdLst
    elements = [deepcopy(sldIdLst[i]) for i in new_order]
    # Remove all existing
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    # Re-add in new order
    for el in elements:
        sldIdLst.append(el)


def replace_text_in_slide(slide, shape_name: str, new_text: str) -> bool:
    """Replace text in a shape identified by name. Preserves formatting of first run."""
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    # Keep formatting from first run, replace text
                    first_run = para.runs[0]
                    first_run.text = new_text
                    # Remove extra runs
                    for run in para.runs[1:]:
                        run.text = ""
                    return True
            # Fallback: set text directly
            shape.text_frame.text = new_text
            return True
    return False


def replace_image_in_slide(slide, shape_name: str, new_image_path: str) -> bool:
    """Replace an image in a shape. Keeps position and size, swaps the image blob."""
    new_path = _resolve_asset_path(new_image_path)

    for shape in slide.shapes:
        if shape.name == shape_name and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Store position and size
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height

            # Remove old shape element
            sp = shape._element
            parent = sp.getparent()
            parent.remove(sp)

            # Add new image with same dimensions
            slide.shapes.add_picture(str(new_path), left, top, width, height)
            return True
    return False


def copy_slide_within_presentation(prs: Presentation, src_index: int) -> int:
    """Duplicate a slide within the same presentation. Returns new slide index."""
    template_slide = prs.slides[src_index]
    slide_layout = template_slide.slide_layout

    # Add new slide with same layout
    new_slide = prs.slides.add_slide(slide_layout)
    new_index = len(prs.slides) - 1

    # Copy shapes
    for shape in template_slide.shapes:
        el = deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # Copy background if set
    try:
        bg = template_slide.background
        if bg._element is not None:
            new_bg = deepcopy(bg._element)
            new_slide.background._element.getparent().replace(
                new_slide.background._element, new_bg
            )
    except Exception:
        pass

    # Remove the default placeholder shapes that add_slide creates
    # (we already copied all shapes from the source)
    spTree = new_slide.shapes._spTree
    for sp in list(spTree):
        # Remove placeholder shapes that were auto-added by the layout
        if sp.tag.endswith("}sp"):
            # Check if this is a duplicate of what we copied
            # Simple heuristic: if shape was in original, keep it
            pass

    return new_index


def copy_slide_cross_presentation(src_prs: Presentation, src_index: int,
                                   dst_prs: Presentation) -> int:
    """Copy a slide from one presentation to another.

    This handles the complex task of copying shapes, images, and relationships.
    Returns the index of the new slide in dst_prs.
    """
    src_slide = src_prs.slides[src_index]

    # Use first available layout in destination
    slide_layout = dst_prs.slide_layouts[0]
    new_slide = dst_prs.slides.add_slide(slide_layout)
    new_index = len(dst_prs.slides) - 1

    # Clear auto-generated shapes from the layout
    spTree = new_slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag.endswith("}sp") or sp.tag.endswith("}pic"):
            spTree.remove(sp)

    # Copy each shape from source
    for shape in src_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # For images, we need to copy the image data
            try:
                image_blob = shape.image.blob
                content_type = shape.image.content_type

                # Determine file extension from content type
                ext_map = {
                    "image/png": ".png",
                    "image/jpeg": ".jpg",
                    "image/gif": ".gif",
                    "image/bmp": ".bmp",
                    "image/svg+xml": ".svg",
                    "image/webp": ".webp",
                }
                ext = ext_map.get(content_type, ".png")

                # Write blob to temp file and add to destination
                import tempfile
                with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                    tmp.write(image_blob)
                    tmp_path = tmp.name

                new_slide.shapes.add_picture(
                    tmp_path,
                    shape.left, shape.top,
                    shape.width, shape.height,
                )
                os.unlink(tmp_path)
            except Exception:
                # Fallback: copy XML element directly
                el = deepcopy(shape._element)
                spTree.append(el)
        else:
            # For non-image shapes, copy the XML element
            el = deepcopy(shape._element)
            spTree.append(el)

    return new_index


# ---------------------------------------------------------------------------
# CLI commands
# ---------------------------------------------------------------------------

def cmd_catalog(args):
    """List available slides from the slide index."""
    index = _load_slide_index()
    slides = index["slides"]

    if args.category:
        slides = [s for s in slides if args.category in s["category"]]
    if args.product:
        slides = [s for s in slides if args.product in s["products"]]

    # Compact output for the agent
    output = []
    for s in slides:
        output.append({
            "id": s["id"],
            "category": s["category"],
            "source": s["source_deck"],
            "products": s["products"],
            "description": s["description"][:120],
            "images": len(s["content_zones"]["images"]),
            "text_zones": len(s["content_zones"]["text"]),
            "has_media": s["has_media"],
            "use_cases": s["use_case_tags"],
        })

    print(json.dumps({"total": len(output), "slides": output}, indent=2))


def cmd_layouts(args):
    """List available slide layouts for creating new slides."""
    # Describe the layouts available in the template presentations
    layouts = [
        {
            "id": "feature-showcase",
            "name": "Feature Showcase",
            "description": "Title + large screenshot + description text",
            "placeholders": ["title (text)", "body (text)", "screenshot (image)"],
        },
        {
            "id": "two-column",
            "name": "Two-Column Comparison",
            "description": "Left/right columns with images and text",
            "placeholders": ["title (text)", "left_text (text)", "right_text (text)",
                             "left_image (image)", "right_image (image)"],
        },
        {
            "id": "stats",
            "name": "Stats / Social Proof",
            "description": "Large number with supporting text",
            "placeholders": ["stat_number (text)", "stat_label (text)", "description (text)"],
        },
        {
            "id": "case-study",
            "name": "Case Study",
            "description": "Client logo + results + screenshot",
            "placeholders": ["client_name (text)", "results (text)", "screenshot (image)"],
        },
        {
            "id": "section-divider",
            "name": "Section Divider",
            "description": "Dark background with section title",
            "placeholders": ["title (text)", "subtitle (text)"],
        },
        {
            "id": "closing",
            "name": "Closing / CTA",
            "description": "Contact info and next steps",
            "placeholders": ["title (text)", "contact (text)"],
        },
    ]
    print(json.dumps({"layouts": layouts}, indent=2))


def cmd_presets(args):
    """List use-case presets."""
    index = _load_slide_index()
    presets = index.get("presets", {})
    output = {}
    for name, preset in presets.items():
        output[name] = {
            "description": preset["description"],
            "product": preset["product"],
            "slide_count": len(preset["slides"]),
            "slides": preset["slides"],
        }
    print(json.dumps({"presets": output}, indent=2))


def cmd_assets(args):
    """List assets from the asset library."""
    manifest_path = ASSETS_DIR / "manifest.json"
    if not manifest_path.exists():
        print(json.dumps({"error": "Asset manifest not found. Run: python3 ppt/extract_assets.py"}))
        return

    with open(manifest_path) as f:
        manifest = json.load(f)

    assets = manifest.get("assets", [])

    if args.category:
        assets = [a for a in assets if args.category in a.get("media_type", "")]

    # Compact output
    output = []
    for a in assets[:100]:  # Limit to first 100
        output.append({
            "id": a.get("id", ""),
            "path": a.get("path", ""),
            "type": a.get("media_type", ""),
            "size": a.get("file_size_bytes", 0),
            "tiny": a.get("is_tiny", False),
        })

    print(json.dumps({
        "total": len(manifest.get("assets", [])),
        "shown": len(output),
        "assets": output,
    }, indent=2))


def cmd_generate(args):
    """Generate a deck from a use-case preset."""
    index = _load_slide_index()
    presets = index.get("presets", {})

    if args.preset not in presets:
        print(json.dumps({
            "error": f"Unknown preset: {args.preset}",
            "available": list(presets.keys()),
        }))
        sys.exit(1)

    preset = presets[args.preset]
    slide_ids = preset["slides"]

    # Build a deck spec from the preset
    spec = {
        "client": args.client,
        "output": args.output or f"Product-deck-{args.client.replace(' ', '-')}-{args.preset}.pptx",
        "slides": [{"action": "use", "slide_id": sid} for sid in slide_ids],
    }

    _execute_compose(spec, index)


def cmd_compose(args):
    """Generate a deck from a JSON deck spec."""
    spec_path = Path(args.spec)
    if not spec_path.exists():
        print(json.dumps({"error": f"Spec file not found: {args.spec}"}))
        sys.exit(1)

    with open(spec_path) as f:
        spec = json.load(f)

    index = _load_slide_index()
    _execute_compose(spec, index)


def _execute_compose(spec: dict, index: dict):
    """Core composition engine. Takes a deck spec and produces a PPTX."""
    slides_spec = spec.get("slides", [])
    client_name = spec.get("client", "Client")
    output_filename = spec.get("output", f"Product-deck-{client_name.replace(' ', '-')}.pptx")

    # Build a lookup of slide_id → slide info
    slide_lookup = {s["id"]: s for s in index["slides"]}

    # Classify each spec entry
    resolved_entries = []  # (spec_order, action, entry, slide_info_or_None)
    for i, entry in enumerate(slides_spec):
        action = entry.get("action", "use")
        if action == "create":
            resolved_entries.append((i, action, entry, None))
        else:
            slide_id = entry.get("slide_id", "")
            if slide_id not in slide_lookup:
                print(json.dumps({
                    "error": f"Unknown slide_id: {slide_id}",
                    "available_ids": sorted(slide_lookup.keys()),
                }))
                sys.exit(1)
            resolved_entries.append((i, action, entry, slide_lookup[slide_id]))

    if not resolved_entries:
        print(json.dumps({"error": "No slides specified in the deck spec"}))
        sys.exit(1)

    # Find the primary source (file with most slides referenced)
    source_counts: dict[str, int] = {}
    for _, action, _, info in resolved_entries:
        if info:
            f = info["source_file"]
            source_counts[f] = source_counts.get(f, 0) + 1
    primary_source = max(source_counts, key=source_counts.get) if source_counts else None

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUTPUT_DIR / output_filename

    if primary_source:
        src_path = PPT_DIR / primary_source
        shutil.copy2(str(src_path), str(output_path))
    else:
        prs = Presentation()
        prs.slide_width = Emu(9144000)
        prs.slide_height = Emu(5143500)
        prs.save(str(output_path))

    prs = Presentation(str(output_path))
    total_slides_original = len(prs.slides)

    # --- Phase 1: Keep only needed primary slides, delete the rest ---
    primary_keep_indices = set()  # 0-based original indices to keep
    for _, action, _, info in resolved_entries:
        if info and info["source_file"] == primary_source:
            idx = info["source_slide_index"] - 1
            if 0 <= idx < total_slides_original:
                primary_keep_indices.add(idx)

    indices_to_delete = sorted(
        [i for i in range(total_slides_original) if i not in primary_keep_indices],
        reverse=True,
    )
    for idx in indices_to_delete:
        try:
            delete_slide(prs, idx)
        except Exception as e:
            print(f"Warning: Could not delete slide {idx}: {e}", file=sys.stderr)

    # After deletion, map original indices → new positions
    remaining_originals = sorted(primary_keep_indices)
    orig_to_new = {orig: new for new, orig in enumerate(remaining_originals)}

    # --- Phase 2: Apply modifications to primary slides ---
    for _, action, entry, info in resolved_entries:
        if action != "modify" or not info or info["source_file"] != primary_source:
            continue
        mods = entry.get("modifications", {})
        new_idx = orig_to_new.get(info["source_slide_index"] - 1)
        if new_idx is None or new_idx >= len(prs.slides):
            continue
        slide = prs.slides[new_idx]
        for text_mod in mods.get("text", []):
            replace_text_in_slide(slide, text_mod.get("shape_name", ""), text_mod.get("value", ""))
        for img_mod in mods.get("images", []):
            asset = img_mod.get("asset", "")
            if asset:
                try:
                    replace_image_in_slide(slide, img_mod.get("shape_name", ""), asset)
                except FileNotFoundError as e:
                    print(f"Warning: {e}", file=sys.stderr)

    # Save and reload to reconcile internal state after deletions.
    # This prevents new slides from conflicting with deleted slide filenames.
    prs.save(str(output_path))
    prs = Presentation(str(output_path))

    # --- Phase 3: Copy cross-deck slides & create new slides (appended at end) ---
    # Track where each spec entry ends up
    spec_to_final_idx: dict[int, int] = {}

    # Map primary slides first
    for spec_order, action, entry, info in resolved_entries:
        if info and info["source_file"] == primary_source:
            new_idx = orig_to_new.get(info["source_slide_index"] - 1)
            if new_idx is not None:
                spec_to_final_idx[spec_order] = new_idx

    # Append cross-deck and created slides
    for spec_order, action, entry, info in resolved_entries:
        if action == "create":
            layout_id = entry.get("layout", "feature-showcase")
            content = entry.get("content", {})
            _create_slide_from_layout(prs, layout_id, content)
            spec_to_final_idx[spec_order] = len(prs.slides) - 1
        elif info and info["source_file"] != primary_source:
            try:
                other_prs = _load_presentation(info["source_file"])
                src_idx = info["source_slide_index"] - 1
                if 0 <= src_idx < len(other_prs.slides):
                    copy_slide_cross_presentation(other_prs, src_idx, prs)
                    spec_to_final_idx[spec_order] = len(prs.slides) - 1
            except Exception as e:
                print(f"Warning: Could not copy slide from {info['source_file']}: {e}",
                      file=sys.stderr)

    # --- Phase 4: Reorder slides to match spec order ---
    # Build the desired ordering: spec_order 0,1,2,... → current indices
    desired_order = []
    for spec_order in range(len(slides_spec)):
        if spec_order in spec_to_final_idx:
            desired_order.append(spec_to_final_idx[spec_order])

    if len(desired_order) == len(prs.slides):
        try:
            reorder_slides(prs, desired_order)
        except Exception as e:
            print(f"Warning: Could not reorder slides: {e}", file=sys.stderr)

    prs.save(str(output_path))

    result = {
        "status": "success",
        "output_path": str(output_path),
        "slide_count": len(prs.slides),
        "client": client_name,
    }
    print(json.dumps(result, indent=2))


def _create_slide_from_layout(prs: Presentation, layout_id: str, content: dict):
    """Create a new slide from a layout template and fill with content."""
    # Use first available layout (or find one by name)
    layout = prs.slide_layouts[0]
    for sl in prs.slide_layouts:
        if layout_id.replace("-", " ").lower() in sl.name.lower():
            layout = sl
            break

    slide = prs.slides.add_slide(layout)

    # Load design tokens for styling
    tokens_path = PPT_DIR / "analysis" / "design_tokens.json"
    if tokens_path.exists():
        with open(tokens_path) as f:
            tokens = json.load(f)
        primary_font = tokens["fonts"]["families"][0]["name"] if tokens["fonts"]["families"] else "Inter"
    else:
        primary_font = "Inter"

    # Add content based on layout type
    title_text = content.get("title", "")
    body_text = content.get("body", "")
    image_path = content.get("image") or content.get("screenshot")

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    if title_text:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(1),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.name = primary_font
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)

    if body_text:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(4.5), Inches(3.5),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = body_text
        run = p.runs[0]
        run.font.size = Pt(14)
        run.font.name = primary_font
        run.font.color.rgb = RGBColor(0x24, 0x24, 0x24)

    if image_path:
        try:
            resolved = _resolve_asset_path(image_path)
            if body_text:
                # Right side of slide
                slide.shapes.add_picture(
                    str(resolved),
                    Inches(5.5), Inches(1.2),
                    Inches(4), Inches(3.5),
                )
            else:
                # Center of slide
                slide.shapes.add_picture(
                    str(resolved),
                    Inches(1), Inches(1.2),
                    Inches(8), Inches(4),
                )
        except FileNotFoundError as e:
            print(f"Warning: {e}", file=sys.stderr)


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="PPTX Presentation Generator — Compose client-specific product decks",
    )
    subparsers = parser.add_subparsers(dest="command")

    # catalog
    cat_parser = subparsers.add_parser("catalog", help="List available slides")
    cat_parser.add_argument("--category", "-c", help="Filter by category substring")
    cat_parser.add_argument("--product", "-p", help="Filter by product (in-app, off-app)")

    # layouts
    subparsers.add_parser("layouts", help="List slide layouts for new slides")

    # presets
    subparsers.add_parser("presets", help="List use-case presets")

    # assets
    assets_parser = subparsers.add_parser("assets", help="List asset library")
    assets_parser.add_argument("--category", "-c", help="Filter by media type (image, video)")

    # generate
    gen_parser = subparsers.add_parser("generate", help="Generate from use-case preset")
    gen_parser.add_argument("--preset", "-p", required=True, help="Preset name")
    gen_parser.add_argument("--client", "-n", required=True, help="Client name")
    gen_parser.add_argument("--output", "-o", help="Output filename")

    # compose
    comp_parser = subparsers.add_parser("compose", help="Generate from JSON deck spec")
    comp_parser.add_argument("--spec", "-s", required=True, help="Path to deck spec JSON")

    args = parser.parse_args()

    if args.command == "catalog":
        cmd_catalog(args)
    elif args.command == "layouts":
        cmd_layouts(args)
    elif args.command == "presets":
        cmd_presets(args)
    elif args.command == "assets":
        cmd_assets(args)
    elif args.command == "generate":
        cmd_generate(args)
    elif args.command == "compose":
        cmd_compose(args)
    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()
